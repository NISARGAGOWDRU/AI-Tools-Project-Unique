import streamlit as st
import google.generativeai as genai
import base64
from pptx import Presentation
from pptx.dml.color import RGBColor  # Import RGBColor
from pptx.util import Inches, Pt
from io import BytesIO

# Configure the Gemini API key (replace with your actual Gemini API key)
API_KEY = "AIzaSyBt3rmdSvRMU8lo3HnjJFELqzK1boA43m0"  # Replace with your actual key
genai.configure(api_key=API_KEY)

def generate_slide_content(topic, slide_number, model_name):
    """Generates content for each slide using the specified model."""
    try:
        model = genai.GenerativeModel(model_name=model_name)

        # Create content for each slide
        content_prompt = f"Provide content for slide {slide_number} about the topic '{topic}'."
        content_response = model.generate_content(content_prompt)
        # Some SDK responses expose text on .text, others on .content; handle both
        content = getattr(content_response, "text", None) or getattr(content_response, "content", None)
        if content is None:
            # Fallback: convert response to string
            content = str(content_response)

        return content.strip()
    except Exception as e:
        st.error(f"An error occurred: {e}")
        # Try to list available models to help the user pick a supported one
        try:
            models = genai.list_models()
            st.info("Available models:")
            for m in models:
                # model objects may be dict-like or have a name attribute
                name = m.get("name") if isinstance(m, dict) else getattr(m, "name", str(m))
                st.write(f"- {name}")
        except Exception:
            st.info("Could not retrieve model list from the API.")
        return None

def create_ppt(content_list):
    """Creates a PowerPoint presentation from the content list with design enhancements."""
    prs = Presentation()

    # Set slide width and height to 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for slide_number, content in enumerate(content_list, start=1):
        slide_layout = prs.slide_layouts[5]  # Blank layout for custom design
        slide = prs.slides.add_slide(slide_layout)

        # Set background color for the slide (light blue)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue color (corrected)

        # Title of the slide
        title_shape = slide.shapes.title
        title_shape.text = f"Slide {slide_number}"

        # Add content box with proper positioning and font size
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(11)
        height = Inches(5)

        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.add_paragraph()
        p.text = content

        # Adjust font size and style
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)  # Set font size to 18pt

    # Save the PowerPoint to a BytesIO object for download
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

def main():
    """Streamlit main function to run the app."""
    st.title("Text to PowerPoint Generator with Gemini API")

    topic = st.text_input("Enter a topic for your presentation:")
    slide_count = st.number_input("Enter the number of slides", min_value=1, max_value=20, value=5)

    # Try to fetch available models and show a dropdown so users pick a supported one
    model_choices = []
    default_model = "models/gemini-2.5-flash"
    try:
        api_models = genai.list_models()
        for m in api_models:
            name = m.get("name") if isinstance(m, dict) else getattr(m, "name", None)
            if name:
                model_choices.append(name)
    except Exception:
        # If listing models fails, fall back to a sensible default list
        model_choices = [
            "models/gemini-2.5-flash",
            "models/gemini-2.5-pro",
            "models/gemini-flash-latest",
            "models/gemini-pro-latest",
        ]

    if default_model not in model_choices:
        model_choices.insert(0, default_model)

    model_name = st.selectbox("Model (pick a supported model)", options=model_choices, index=0)

    if topic and slide_count:
        st.write(f"Generating {slide_count} slides for the topic: {topic}...")

        # Generate slide content
        content_list = []
        for slide_number in range(1, slide_count + 1):
            content = generate_slide_content(topic, slide_number, model_name)
            if content:
                content_list.append(content)

        # Create the PowerPoint presentation
        if content_list:
            ppt_io = create_ppt(content_list)

            # Provide a download link for the PowerPoint file
            st.subheader("Download your PowerPoint presentation:")
            st.download_button(
                label="Download PPT",
                data=ppt_io,
                file_name=f"{topic}_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

if __name__ == "__main__":
    main()
